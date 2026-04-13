"""
Ingest — the mouth of the fish.

Reads files, extracts text, chunks by natural boundaries.
Each file is an "exchange" — R(n) grows with each one ingested.

Supported formats (stdlib only unless marked):
  .md, .markdown    — markdown, chunked by headers
  .txt, .log        — plain text, chunked by paragraphs
  .rst, .tex, .org  — writing formats, chunked as text
  .html, .htm       — HTML, tags stripped
  .csv, .tsv        — tabular data, formatted as text rows
  .json             — JSON, pretty-printed
  .jsonl, .ndjson   — line-delimited JSON, one chunk per record
  .yaml, .yml       — YAML, via pyyaml if installed, else read as text
  .toml, .ini, .cfg — config files, read as text
  .xml              — XML, tags stripped
  .py, .js, .ts, .go, .rs — source code, chunked by def/class/fn
  .pdf              — PDF, via PyMuPDF or pdfplumber (optional deps)
  .docx             — Word, via python-docx (optional dep)
  .pptx             — PowerPoint, via python-pptx (optional dep)
  .rtf              — Rich Text, via striprtf (optional) or regex strip

Unknown suffixes fall through to read_text, so a directory of mixed
content will still be eaten. Set strict=True on ingest_directory to
keep only extensions in READERS.

The chunking is semantic, not mechanical. A 512-token window doesn't
know where meaning lives. Headers, paragraphs, and topic shifts do.
"""

from __future__ import annotations


import re
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional


@dataclass
class Chunk:
    """A unit of ingested content."""
    text: str
    source: str           # file path
    section: str = ""     # header/section name if available
    chunk_type: str = ""  # "narrative", "data", "code", "metadata"
    position: int = 0     # order within source file


def read_markdown(path: Path) -> list[Chunk]:
    """Read a markdown file, chunk by headers."""
    text = path.read_text(encoding="utf-8", errors="replace")
    chunks = []
    current_section = ""
    current_lines = []
    position = 0

    for line in text.split("\n"):
        if line.startswith("#"):
            # Save previous section
            if current_lines:
                body = "\n".join(current_lines).strip()
                if body and len(body) > 20:
                    chunks.append(Chunk(
                        text=body,
                        source=str(path),
                        section=current_section,
                        chunk_type="narrative",
                        position=position,
                    ))
                    position += 1
            current_section = line.lstrip("#").strip()
            current_lines = [line]
        else:
            current_lines.append(line)

    # Last section
    if current_lines:
        body = "\n".join(current_lines).strip()
        if body and len(body) > 20:
            chunks.append(Chunk(
                text=body,
                source=str(path),
                section=current_section,
                chunk_type="narrative",
                position=position,
            ))

    # If no headers found, chunk by paragraphs
    if not chunks:
        chunks = chunk_by_paragraphs(text, str(path))

    return chunks


def read_text(path: Path) -> list[Chunk]:
    """Read a plain text file, chunk by paragraphs."""
    text = path.read_text(encoding="utf-8", errors="replace")
    return chunk_by_paragraphs(text, str(path))


def read_pdf(path: Path) -> list[Chunk]:
    """Read a PDF file. Uses PyMuPDF if available, falls back to basic extraction."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(str(path))
        chunks = []
        for page_num, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                # Try to split by apparent sections
                sections = re.split(r'\n(?=[A-Z][A-Z\s]{5,})', text)
                for i, section in enumerate(sections):
                    section = section.strip()
                    if section and len(section) > 30:
                        chunks.append(Chunk(
                            text=section,
                            source=str(path),
                            section=f"page_{page_num + 1}",
                            chunk_type="narrative",
                            position=page_num * 100 + i,
                        ))
        doc.close()
        return chunks
    except ImportError:
        # Fallback: try pdfplumber
        try:
            import pdfplumber
            chunks = []
            with pdfplumber.open(str(path)) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    if text and text.strip() and len(text.strip()) > 30:
                        chunks.append(Chunk(
                            text=text.strip(),
                            source=str(path),
                            section=f"page_{page_num + 1}",
                            chunk_type="narrative",
                            position=page_num,
                        ))
            return chunks
        except ImportError:
            print(f"  [skip] No PDF reader available for {path.name}")
            return []


def read_docx(path: Path) -> list[Chunk]:
    """Read a Word document."""
    try:
        from docx import Document
        doc = Document(str(path))
        chunks = []
        current_section = ""
        current_lines = []
        position = 0

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Detect headers by style
            if para.style and para.style.name and "Heading" in para.style.name:
                if current_lines:
                    body = "\n".join(current_lines)
                    if len(body) > 20:
                        chunks.append(Chunk(
                            text=body,
                            source=str(path),
                            section=current_section,
                            chunk_type="narrative",
                            position=position,
                        ))
                        position += 1
                current_section = text
                current_lines = []
            else:
                current_lines.append(text)

        if current_lines:
            body = "\n".join(current_lines)
            if len(body) > 20:
                chunks.append(Chunk(
                    text=body,
                    source=str(path),
                    section=current_section,
                    chunk_type="narrative",
                    position=position,
                ))

        if not chunks:
            # No structure found, dump as one chunk
            all_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            if all_text and len(all_text) > 20:
                chunks.append(Chunk(
                    text=all_text,
                    source=str(path),
                    section="",
                    chunk_type="narrative",
                    position=0,
                ))

        return chunks
    except ImportError:
        print(f"  [skip] python-docx not installed for {path.name}")
        return []


def read_json(path: Path) -> list[Chunk]:
    """Read a JSON file — tool definitions, configs, schemas."""
    import json
    text = path.read_text(encoding="utf-8", errors="replace")
    try:
        data = json.loads(text)
        # Pretty-print for readability
        formatted = json.dumps(data, indent=2)
        return [Chunk(
            text=formatted,
            source=str(path),
            section="",
            chunk_type="data",
            position=0,
        )]
    except json.JSONDecodeError:
        return [Chunk(text=text, source=str(path), chunk_type="data")]


def read_jsonl(path: Path) -> list[Chunk]:
    """Read line-delimited JSON — one chunk per record."""
    import json
    chunks = []
    with path.open("r", encoding="utf-8", errors="replace") as f:
        for i, line in enumerate(f):
            line = line.strip()
            if not line:
                continue
            try:
                data = json.loads(line)
                text = json.dumps(data, indent=2)
            except json.JSONDecodeError:
                text = line
            if len(text) > 20:
                chunks.append(Chunk(
                    text=text,
                    source=str(path),
                    section=f"record_{i}",
                    chunk_type="data",
                    position=i,
                ))
    return chunks


def read_html(path: Path) -> list[Chunk]:
    """Read an HTML file, strip tags, chunk by paragraphs.

    Uses html.parser from stdlib — no extra deps. Preserves link text and
    headings by wrapping them in plain prose.
    """
    from html.parser import HTMLParser

    class _Stripper(HTMLParser):
        def __init__(self):
            super().__init__()
            self.parts: list[str] = []
            self.in_script = False
            self.in_style = False

        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style"):
                if tag == "script":
                    self.in_script = True
                else:
                    self.in_style = True
            elif tag in ("p", "br", "div", "li", "tr"):
                self.parts.append("\n")
            elif tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
                self.parts.append("\n\n# ")

        def handle_endtag(self, tag):
            if tag == "script":
                self.in_script = False
            elif tag == "style":
                self.in_style = False
            elif tag in ("h1", "h2", "h3", "h4", "h5", "h6", "p"):
                self.parts.append("\n")

        def handle_data(self, data):
            if not (self.in_script or self.in_style):
                self.parts.append(data)

        def text(self) -> str:
            return "".join(self.parts)

    text = path.read_text(encoding="utf-8", errors="replace")
    parser = _Stripper()
    try:
        parser.feed(text)
        cleaned = re.sub(r"[ \t]+", " ", parser.text())
        cleaned = re.sub(r"\n\s*\n\s*\n+", "\n\n", cleaned).strip()
    except Exception:
        cleaned = re.sub(r"<[^>]+>", " ", text)
    return chunk_by_paragraphs(cleaned, str(path))


def read_csv(path: Path) -> list[Chunk]:
    """Read CSV/TSV, format rows as readable text. Header-aware."""
    import csv
    chunks = []
    sep = "\t" if path.suffix.lower() == ".tsv" else ","
    try:
        with path.open("r", encoding="utf-8", errors="replace", newline="") as f:
            reader = csv.reader(f, delimiter=sep)
            rows = list(reader)
    except Exception:
        return [Chunk(text=path.read_text(encoding="utf-8", errors="replace"),
                      source=str(path), chunk_type="data")]
    if not rows:
        return []
    header = rows[0] if rows else []
    # Emit a summary chunk + one chunk per row when rows are long enough
    summary = f"CSV: {path.name}\nColumns: {', '.join(header)}\nRows: {len(rows) - 1}"
    chunks.append(Chunk(text=summary, source=str(path), section="header",
                        chunk_type="data", position=0))
    for i, row in enumerate(rows[1:], start=1):
        if not any(c.strip() for c in row):
            continue
        if header and len(header) == len(row):
            pairs = [f"{h}: {v}" for h, v in zip(header, row) if v]
            text = "\n".join(pairs)
        else:
            text = sep.join(row)
        if len(text) > 20:
            chunks.append(Chunk(text=text, source=str(path),
                                section=f"row_{i}", chunk_type="data", position=i))
    return chunks


def read_yaml(path: Path) -> list[Chunk]:
    """Read YAML via pyyaml if installed, else treat as text."""
    text = path.read_text(encoding="utf-8", errors="replace")
    try:
        import yaml
        data = yaml.safe_load(text)
        import json
        formatted = json.dumps(data, indent=2, default=str)
        return [Chunk(text=formatted, source=str(path), chunk_type="data", position=0)]
    except ImportError:
        return [Chunk(text=text, source=str(path), chunk_type="data", position=0)]
    except Exception:
        return [Chunk(text=text, source=str(path), chunk_type="data", position=0)]


def read_pptx(path: Path) -> list[Chunk]:
    """Read a PowerPoint. Optional dep: python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        print(f"  [skip] python-pptx not installed for {path.name}")
        return []
    try:
        prs = Presentation(str(path))
    except Exception as e:
        print(f"  [skip] failed to open {path.name}: {e}")
        return []
    chunks = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
        body = "\n".join(p for p in parts if p)
        if body and len(body) > 20:
            chunks.append(Chunk(
                text=body,
                source=str(path),
                section=f"slide_{i}",
                chunk_type="narrative",
                position=i,
            ))
    return chunks


def read_rtf(path: Path) -> list[Chunk]:
    """Read RTF. Optional dep: striprtf, else regex strip."""
    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        from striprtf.striprtf import rtf_to_text
        text = rtf_to_text(raw)
    except ImportError:
        # Minimal regex strip — removes control words + braces. Lossy but workable.
        text = re.sub(r"\\[a-z]+-?\d*\s?", "", raw)
        text = re.sub(r"[{}]", "", text)
    except Exception:
        text = raw
    return chunk_by_paragraphs(text, str(path))


def read_xml(path: Path) -> list[Chunk]:
    """Read XML, strip tags."""
    text = path.read_text(encoding="utf-8", errors="replace")
    cleaned = re.sub(r"<[^>]+>", " ", text)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    return chunk_by_paragraphs(cleaned, str(path))


def read_python(path: Path) -> list[Chunk]:
    """Read a Python file, chunk by class/function definitions."""
    text = path.read_text(encoding="utf-8", errors="replace")
    chunks = []
    current_block = []
    current_name = path.name
    position = 0

    for line in text.split("\n"):
        if re.match(r'^(class |def |async def )', line):
            if current_block:
                body = "\n".join(current_block)
                if len(body) > 30:
                    chunks.append(Chunk(
                        text=body,
                        source=str(path),
                        section=current_name,
                        chunk_type="code",
                        position=position,
                    ))
                    position += 1
            current_name = line.split("(")[0].replace("class ", "").replace("def ", "").replace("async def ", "").strip()
            current_block = [line]
        else:
            current_block.append(line)

    if current_block:
        body = "\n".join(current_block)
        if len(body) > 30:
            chunks.append(Chunk(
                text=body,
                source=str(path),
                section=current_name,
                chunk_type="code",
                position=position,
            ))

    return chunks


def chunk_by_paragraphs(text: str, source: str) -> list[Chunk]:
    """Split text into paragraph-level chunks."""
    paragraphs = re.split(r'\n\s*\n', text)
    chunks = []
    for i, para in enumerate(paragraphs):
        para = para.strip()
        if para and len(para) > 20:
            chunks.append(Chunk(
                text=para,
                source=source,
                section="",
                chunk_type="narrative",
                position=i,
            ))
    return chunks


# Reader dispatch
READERS = {
    # Writing formats — chunked by structure
    ".md": read_markdown,
    ".markdown": read_markdown,
    ".txt": read_text,
    ".log": read_text,
    ".rst": read_text,
    ".tex": read_text,
    ".org": read_text,
    # Web / markup
    ".html": read_html,
    ".htm": read_html,
    ".xml": read_xml,
    # Data / config
    ".json": read_json,
    ".jsonl": read_jsonl,
    ".ndjson": read_jsonl,
    ".csv": read_csv,
    ".tsv": read_csv,
    ".yaml": read_yaml,
    ".yml": read_yaml,
    ".toml": read_text,
    ".ini": read_text,
    ".cfg": read_text,
    ".conf": read_text,
    ".env": read_text,
    # Documents
    ".pdf": read_pdf,
    ".docx": read_docx,
    ".pptx": read_pptx,
    ".rtf": read_rtf,
    # Source code — chunked by def/class
    ".py": read_python,
    ".js": read_python,
    ".ts": read_python,
    ".go": read_python,
    ".rs": read_python,
    ".java": read_python,
    ".kt": read_python,
    ".rb": read_python,
    ".php": read_python,
    ".sh": read_text,
    ".bash": read_text,
    ".zsh": read_text,
    ".ps1": read_text,
}

# Extensions we explicitly never try to read (binary formats without a
# text extraction path in this module). A directory walk filters these
# out rather than attempting read_text on them. Aggressive by design —
# better to miss a .bak that happened to be text than ingest a .bin and
# produce garbage crystals.
BINARY_SKIP = {
    # Images
    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif",
    ".ico", ".cur", ".svg", ".psd", ".ai", ".eps", ".indd", ".heic", ".heif",
    ".raw", ".arw", ".cr2", ".nef", ".orf", ".rw2",
    # Audio
    ".mp3", ".wav", ".m4a", ".flac", ".ogg", ".oga", ".aac", ".opus",
    ".wma", ".amr", ".aiff", ".au",
    # Video
    ".mp4", ".mov", ".avi", ".mkv", ".webm", ".wmv", ".flv", ".m4v",
    ".3gp", ".mpg", ".mpeg", ".vob",
    # Archives / installers
    ".zip", ".tar", ".gz", ".bz2", ".xz", ".zst", ".lz4",
    ".7z", ".rar", ".iso", ".dmg", ".pkg", ".msi", ".deb", ".rpm",
    ".apk", ".ipa", ".jar", ".war", ".ear",
    # Executables / libraries / object files
    ".exe", ".dll", ".so", ".dylib", ".class", ".o", ".a", ".lib",
    ".sys", ".bin", ".obj", ".com", ".app",
    # Python + other bytecode
    ".pyc", ".pyo", ".pyd",
    # Databases
    ".db", ".sqlite", ".sqlite3", ".db-journal", ".mdb", ".accdb",
    ".realm", ".leveldb", ".idx", ".pack",
    # Fonts
    ".ttf", ".otf", ".woff", ".woff2", ".eot", ".fon", ".fnt",
    # Backups / editor state / OS cruft that's almost never content
    ".bak", ".old", ".orig", ".tmp", ".temp", ".swp", ".swo", ".swn",
    ".ds_store", ".lnk",
    # Git / VCS internals (directories usually, but just in case)
    ".pack", ".idx",
    # Proprietary office binaries we can't read (.doc/.xls/.ppt — legacy)
    ".doc", ".xls", ".ppt", ".dot", ".xlt", ".pot",
    # Large media containers
    ".psdx", ".aep", ".prproj",
    # Compiled docs
    ".chm", ".mobi", ".azw", ".azw3",
}

# Files larger than this that fall through to read_text get skipped — a
# 50 MB unknown suffix is almost certainly binary-ish and reading it as
# UTF-8 would produce noise and block the ingest for minutes. Known
# readers (pdf/docx/pptx) still handle large files via their own paths.
FALLTHROUGH_MAX_BYTES = 5 * 1024 * 1024  # 5 MB


def read_file_as_text(path: Path) -> str:
    """Read a file through the proper reader and return its joined text.

    This is the "one-shot" shape used by `linafish go` and other callers
    that want a single string per file instead of chunked Chunk objects.
    It still goes through the READERS dispatch, so HTML gets its tags
    stripped, CSV becomes readable rows, YAML pretty-prints, etc.

    Returns empty string if the file can't be read or is binary-skipped.
    Never raises.
    """
    chunks = ingest_file(path)
    if not chunks:
        return ""
    return "\n\n".join(c.text for c in chunks if c.text)


def ingest_file(path: Path) -> list[Chunk]:
    """Ingest a single file. Unknown suffixes fall back to read_text with
    a size guard so we don't try to UTF-8 decode a 500 MB binary blob."""
    suffix = path.suffix.lower()
    if suffix in BINARY_SKIP:
        return []
    # Also skip hidden files (start with dot) and common no-suffix cruft
    if path.name.startswith(".") and suffix == "":
        return []
    reader = READERS.get(suffix)
    if reader is None:
        # Fall through: treat anything we don't explicitly skip as text.
        # But only if the file is small enough to plausibly be text — a
        # giant unknown file is almost certainly binary we missed in
        # BINARY_SKIP, and reading it as UTF-8 pollutes the fish.
        try:
            size = path.stat().st_size
        except Exception:
            return []
        if size > FALLTHROUGH_MAX_BYTES:
            print(f"  [skip] {path.name} ({size:,} bytes) — unknown suffix, too large to fall through")
            return []
        try:
            return read_text(path)
        except Exception:
            return []
    try:
        return reader(path)
    except Exception as e:
        # Never let a single bad file break a directory walk.
        print(f"  [reader error] {path.name}: {type(e).__name__}: {e}")
        return []


def ingest_directory(
    directory: Path,
    extensions: Optional[set[str]] = None,
    recursive: bool = True,
    strict: bool = False,
) -> list[Chunk]:
    """Ingest all files in a directory.

    Each file is an exchange. R(n) grows.

    strict=False (default): any file not in BINARY_SKIP is eaten,
        with unknown extensions falling through to read_text.
    strict=True: only files with suffixes in `extensions` (or READERS
        if extensions is None) are eaten. Legacy behavior.
    """
    all_chunks = []
    pattern = "**/*" if recursive else "*"

    files = sorted(directory.glob(pattern))
    files = [f for f in files if f.is_file()]

    if strict:
        allowed = extensions if extensions is not None else set(READERS.keys())
        files = [f for f in files if f.suffix.lower() in allowed]
    else:
        if extensions is not None:
            files = [f for f in files if f.suffix.lower() in extensions]
        else:
            files = [f for f in files if f.suffix.lower() not in BINARY_SKIP]

    print(f"Ingesting {len(files)} files from {directory}")

    for i, file_path in enumerate(files, 1):
        print(f"  [{i}/{len(files)}] {file_path.name}", end="")
        try:
            chunks = ingest_file(file_path)
            print(f" -> {len(chunks)} chunks")
            all_chunks.extend(chunks)
        except Exception as e:
            print(f" -> ERROR: {e}")

    print(f"Total: {len(all_chunks)} chunks from {len(files)} files")
    return all_chunks
